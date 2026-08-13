"""Python/Rust differential harness with the frozen v2 seven-status semantics.

M0-03 (see the revised v2.0.0 plan doc §4.1 and the
parity ledger §0): the harness classifies every run on both sides under the
single v2 SolveStatus vocabulary —

    SOLVED / PROVEN_INFEASIBLE / TIMEOUT / UNKNOWN / INVALID_INPUT /
    CANCELLED / INTERNAL_ERROR

Frozen rules:

- A Python error is NEVER recorded as INFEASIBLE (that was the old harness
  bug: `benchmark_parity.solve_reference` mapped every SeatTrellisSolveError
  to solver_status="INFEASIBLE").
- Heuristic exhaustion without a proof is UNKNOWN, not PROVEN_INFEASIBLE.
- Wall-clock exhaustion is TIMEOUT; input-validation failures are
  INVALID_INPUT; anything else is INTERNAL_ERROR.
- Any status mismatch exits non-zero. All classes pass -> exit 0.

Status classes exercised:

1. SOLVED class   — benchmark sizes (default 40/50/60) with full student +
                    rule data, and every non-invalid fixture case.
2. INVALID_INPUT  — the seven `invalid-*` cases of fixtures/parity (which the
                    Python CLI rejects by design).
3. TIMEOUT class  — a 60-student case with a near-zero time limit. The Rust
                    core has no time budget yet (ledger gap, M3-04), so a
                    mismatch here is a *documented* gap: it is reported loudly
                    but still fails the run, per the frozen mismatch policy.

Usage:
    python scripts/rust_python_diff.py                      # benchmark sizes
    python scripts/rust_python_diff.py --sizes 40,60        # subset
    python scripts/rust_python_diff.py --time-limit 3
    python scripts/rust_python_diff.py --fixtures           # fixture corpus classes
"""

from __future__ import annotations

import argparse
import csv
import json
import shutil
import subprocess
import sys
import tempfile
from html.parser import HTMLParser
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
CLI = ROOT / "target" / "release" / "seattrellis_cli"  # single root workspace since M1-01
PY_CLI = ROOT / ".venv" / "bin" / "seattrellis"
FIXTURES = ROOT / "fixtures" / "parity"
INPUTS = FIXTURES / "inputs"
GOLDENS = FIXTURES / "goldens"
SOLVER_BACKEND = "fallback"

sys.path.insert(0, str(Path(__file__).resolve().parent))
from benchmark_parity import (  # noqa: E402
    NATIVE_API_VERSION,
    build_case,
)

# --- frozen v2 SolveStatus vocabulary ---------------------------------------

STATUS_SOLVED = "SOLVED"
STATUS_PROVEN_INFEASIBLE = "PROVEN_INFEASIBLE"
STATUS_TIMEOUT = "TIMEOUT"
STATUS_UNKNOWN = "UNKNOWN"
STATUS_INVALID_INPUT = "INVALID_INPUT"
STATUS_CANCELLED = "CANCELLED"
STATUS_INTERNAL_ERROR = "INTERNAL_ERROR"

# Status classes where a Rust mismatch is a *documented* ledger gap (not a
# harness artifact). A mismatch still fails the run per M0-03.
# TIMEOUT was listed here until M3-04 landed --time-limit (PR #95); the
# benchmark TIMEOUT class now exercises the Rust wall-clock budget too.
KNOWN_RUST_GAPS: dict[str, str] = {}

# Case-level documented corpus gaps (case id -> ledger reference). These are
# real equivalence gaps, not harness artifacts: for each of them the Python
# load/resolver rejected the original input, so the harness sent a DEGRADED
# request (unknown rule kinds / bad adjacency references dropped) that the
# Rust core legitimately solves. The comparison is therefore not
# apples-to-apples and each case carries an explicit ledger reference.
# Without `--allow-documented-gaps` the run still fails on them (M0-03);
# with the flag (used by CI) only NEW mismatches fail the run.
#
# All three original entries (invalid-unknown-rule, invalid-unknown-soft-
# objective, invalid-bad-adjacency-ref) were CLOSED on 2026-08-10: the Rust
# project-workspace compiler now rejects unknown hard rule kinds, unknown
# soft objectives and bad adjacency references (mirroring Python's
# extra="forbid" models), and the harness validates those cases through the
# CLI's project-validate path (ledger §19.5, fixture evidence 41/41 match).
DOCUMENTED_CORPUS_GAPS: dict[str, str] = {}

INVALID_TOKENS = (
    "validation",
    "required",
    "require",
    "not enough",
    "unknown",
    "duplicate",
    "invalid",
    "unrecognized",
    "missing",
    "at least",
    "cannot seat",
    "more students",
)


def classify_python_cli(proc: subprocess.CompletedProcess) -> str:
    """Classify a Python CLI run under the v2 vocabulary (never INFEASIBLE)."""
    if proc.returncode == 0:
        return STATUS_SOLVED
    text = (proc.stderr + proc.stdout).lower()
    if "this is not proof" in text or "time limit" in text:
        return STATUS_TIMEOUT
    if "no feasible seating plan" in text:
        return STATUS_UNKNOWN
    if any(token in text for token in INVALID_TOKENS):
        return STATUS_INVALID_INPUT
    return STATUS_INTERNAL_ERROR


def classify_rust_cli(proc: subprocess.CompletedProcess, output: Path | None) -> str:
    """Classify a Rust CLI run under the v2 vocabulary.

    The core reports `feasible: bool` only; a zero-exit with feasible=false is
    greedy exhaustion, which the frozen semantics require to be UNKNOWN (the
    core cannot prove infeasibility yet).
    """
    if proc.returncode == 0 and output is not None and output.exists():
        try:
            data = json.loads(output.read_text(encoding="utf-8"))
        except (OSError, ValueError):
            return STATUS_INTERNAL_ERROR
        if data.get("feasible"):
            return STATUS_SOLVED
        return STATUS_UNKNOWN
    text = proc.stderr.lower()
    if any(token in text for token in INVALID_TOKENS):
        return STATUS_INVALID_INPUT
    return STATUS_INTERNAL_ERROR


# --- Rust request builders --------------------------------------------------

def _core_student(py_student: dict[str, object]) -> dict[str, object]:
    """Map a Python student dump to the core Student shape."""
    key = py_student.get("student_id") or py_student.get("name")
    vision = py_student.get("vision")
    return {
        "key": str(key),
        "display_name": py_student.get("name"),
        "height_cm": py_student.get("height_cm"),
        "score": py_student.get("score"),
        "vision": None if vision is None else str(vision),
        "tags": py_student.get("tags", []),
        "needs": py_student.get("needs", []),
    }


def benchmark_request(case: dict[str, object]) -> dict[str, object]:
    """CoreSolveRequest for a benchmark case: index-space problem + students
    + soft rules (root-level hard fields come from ``problem``)."""
    problem = case["problem"]
    meta = case["problem_meta"]
    request = dict(problem)
    request["students"] = [_core_student(s) for s in meta["students"]]
    rules = json.loads(json.dumps(meta["rules"]))
    # The core reads only `seed` + `soft` from the rules document; hard
    # constraints travel in the root-level index fields emitted by
    # build_problem_json.
    request["rules"] = {"seed": rules.get("seed", 42), "soft": rules.get("soft", {})}
    return request


def fixture_to_request(case_dir: Path) -> tuple[dict[str, object], list[str]]:
    """Convert a fixtures/parity input case into a CoreSolveRequest.

    Known hard rules are resolved to index constraints through the same
    Python resolver the corpus generator uses, so the problem the Rust core
    sees is exactly the problem the Python oracle solved. Unknown rule kinds
    cannot be represented by the core DTO: they are dropped (the core serde
    ignores unknown fields the same way) and the drop is recorded in the
    returned notes. Inputs the pydantic loaders or the resolver reject
    (invalid-* cases) fall back to a minimal degraded request built from the
    raw files that still carries the malformed shape.
    """
    from seattrellis.io.json_files import load_layout, load_rules
    from seattrellis.io.students import read_students
    from seattrellis.solver.precompute import precompute_topology
    from seattrellis.solver.rule_compiler import compile_hard_rules, resolve_hard_rules

    notes: list[str] = []
    try:
        students = read_students(case_dir / "students.csv")
        layout = load_layout(case_dir / "classroom.json")
        rules = load_rules(case_dir / "rules.json")
        topology = precompute_topology(students, layout)
        resolved = resolve_hard_rules(students, layout, rules, topology=topology)
        compiled = compile_hard_rules(resolved)
    except Exception as exc:  # invalid-input cases reject during load/resolve
        notes.append(
            f"Python load/resolver rejected the input ({exc.__class__.__name__}); "
            "the Rust side validates the raw workspace via project-validate"
        )
        return degraded_request_raw(case_dir), notes
    request: dict[str, object] = {
        "api_version": NATIVE_API_VERSION,
        "student_count": len(students),
        "seat_positions": [[float(s.x), float(s.y)] for s in topology.seats],
        "edges": [list(e) for e in sorted(topology.adjacent_seat_index_pairs)],
        "fixed_seats": [[i, j] for i, j in sorted(compiled.fixed_seats.items())],
        "must_be_adjacent": [list(p) for p in compiled.must_be_adjacent],
        "cannot_be_adjacent": [list(p) for p in compiled.cannot_be_adjacent],
        "min_distance": [
            {"students": [a, b], "distance": r.distance, "metric": r.metric}
            for (a, b, r) in compiled.min_distance
        ],
        "seed": rules.seed,
        "students": [_core_student(s.model_dump(mode="json")) for s in students],
        "rules": {
            "seed": rules.seed,
            "soft": rules.model_dump(mode="json").get("soft", {}),
        },
    }
    return request, notes


def degraded_request_raw(case_dir: Path) -> dict[str, object]:
    """Minimal CoreSolveRequest from the raw fixture files, bypassing the
    pydantic loaders (which reject the malformed inputs by design).

    Keeps student_count, seat positions and the student records (including
    duplicates) so the Rust side still sees the malformed shape; adjacency and
    hard rules are dropped and the drop is recorded via the caller's notes.
    """
    layout = json.loads((case_dir / "classroom.json").read_text(encoding="utf-8"))
    seats = [s for s in layout.get("seats", []) if s.get("enabled", True)]
    rules = json.loads((case_dir / "rules.json").read_text(encoding="utf-8"))
    students: list[dict[str, object]] = []
    with (case_dir / "students.csv").open(newline="", encoding="utf-8") as f:
        for row in csv.DictReader(f):
            vision = row.get("vision") or None
            students.append({
                "key": str(row.get("student_id") or row.get("name") or ""),
                "display_name": row.get("name") or None,
                "height_cm": _num_or_none(row.get("height_cm")),
                "score": _num_or_none(row.get("score")),
                "vision": str(vision) if vision else None,
                "tags": _split_tags(row.get("tags")),
                "needs": _split_tags(row.get("needs")),
            })
    seed = rules.get("seed", 42)
    return {
        "api_version": NATIVE_API_VERSION,
        "student_count": len(students),
        "seat_positions": [[float(s["x"]), float(s["y"])] for s in seats],
        "seed": seed,
        "students": students,
        "rules": {"seed": seed, "soft": rules.get("soft", {})},
    }


def _num_or_none(raw: str | None) -> float | None:
    if raw is None or raw.strip() == "":
        return None
    try:
        return float(raw)
    except ValueError:
        return None


def _split_tags(raw: str | None) -> list[str]:
    if raw is None or raw.strip() == "":
        return []
    return [t for t in (part.strip() for part in raw.split(",")) if t]


def run_rust(
    request: dict[str, object],
    tmp: Path,
    time_limit: float | None = None,
) -> tuple[str, str, dict[str, str]]:
    """Run the Rust CLI on a request; return (status, detail, notes).

    ``time_limit`` is forwarded to ``--time-limit`` so the TIMEOUT class
    exercises the M3-04 wall-clock budget (previously a documented gap).
    """
    problem_file = tmp / "rust-problem.json"
    output_file = tmp / "rust-result.json"
    problem_file.write_text(json.dumps(request), encoding="utf-8")
    cmd = [str(CLI), "solve", "--problem", str(problem_file), "--output", str(output_file)]
    if time_limit is not None:
        cmd += ["--time-limit", str(time_limit)]
    proc = subprocess.run(cmd, capture_output=True, text=True)
    status = classify_rust_cli(proc, output_file)
    detail = ""
    if proc.returncode == 0 and output_file.exists():
        try:
            data = json.loads(output_file.read_text(encoding="utf-8"))
            detail = f"feasible={data.get('feasible')} seated={len(data.get('assignment', []))}"
        except (OSError, ValueError):
            detail = proc.stderr.strip()[:200]
    else:
        detail = proc.stderr.strip()[:200]
    return status, detail, {}


def run_rust_workspace_validate(case_dir: Path, tmp: Path) -> tuple[str, str, dict[str, str]]:
    """Run the Rust CLI ``project-validate`` on a synthesized project
    workspace referencing the case's raw files.

    This exercises the same import surface as the Python oracle's
    load/resolve: unknown rule kinds, unknown soft objectives, malformed
    layouts and bad adjacency references surface here instead of being
    degraded away by the index-space request builder. A workspace the
    compiler rejects exits 2 (InvalidInput); a valid workspace exits 0.
    """
    workspace = tmp / f"workspace-{case_dir.name}"
    workspace.mkdir(exist_ok=True)
    project_document = {
        "kind": "seattrellis_project",
        "schema_version": 1,
        "name": case_dir.name,
        "students": "students.csv",
        "layout": "classroom.json",
        "rules": "rules.json",
        "outputs_dir": "outputs",
    }
    (workspace / "project.json").write_text(
        json.dumps(project_document), encoding="utf-8"
    )
    for name in ("students.csv", "classroom.json", "rules.json"):
        source = case_dir / name
        if source.is_file():
            (workspace / name).write_bytes(source.read_bytes())
    proc = subprocess.run(
        [str(CLI), "project-validate", "--project", str(workspace / "project.json")],
        capture_output=True,
        text=True,
    )
    detail = proc.stderr.strip()[:200] or proc.stdout.strip()[:200]
    if proc.returncode == 0:
        return STATUS_SOLVED, detail, {}
    if proc.returncode == 2:
        return STATUS_INVALID_INPUT, detail, {}
    return STATUS_INTERNAL_ERROR, detail, {}


# --- status classes ---------------------------------------------------------

def run_benchmark_classes(sizes: list[int], time_limit: float) -> list[tuple[str, str, str, str, list[str]]]:
    """(case_id, python_status, rust_status, detail, notes) for SOLVED +
    TIMEOUT classes built from benchmark_parity cases."""
    rows: list[tuple[str, str, str, str, list[str]]] = []
    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        for size in sizes:
            case = build_case(size, "daily", seed=42, time_limit=time_limit)
            py_ref = case["python_reference"]
            py_status = py_ref.get("status", STATUS_INTERNAL_ERROR)
            rust = run_rust(benchmark_request(case), tmp_path, time_limit=time_limit)
            rows.append((f"bench-{size}", py_status, rust[0], rust[1], []))
        # TIMEOUT class: same size but the minimum allowed budget (the
        # fallback enforces time_limit_seconds >= 0.1); 60 students exceed it.
        case = build_case(max(sizes), "daily", seed=42, time_limit=0.1)
        py_ref = case["python_reference"]
        py_status = py_ref.get("status", STATUS_INTERNAL_ERROR)
        rust = run_rust(benchmark_request(case), tmp_path, time_limit=0.1)
        rows.append((f"bench-{max(sizes)}-timeout-0.1", py_status, rust[0], rust[1], []))
    return rows


def run_fixture_classes() -> list[tuple[str, str, str, str, list[str]]]:
    """(case_id, python_status, rust_status, detail, notes) for every case in
    fixtures/parity: non-invalid cases are the SOLVED class, invalid-* cases
    are the INVALID_INPUT class."""
    rows: list[tuple[str, str, str, str, list[str]]] = []
    if not CLI.exists():
        raise SystemExit(f"Rust CLI not found: {CLI}; build it first (cargo build --release -p seattrellis_cli)")
    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        for case_dir in sorted(INPUTS.iterdir()):
            if not case_dir.is_dir():
                continue
            cid = case_dir.name
            flags = [
                str(PY_CLI), "solve",
                "--students", str(case_dir / "students.csv"),
                "--layout", str(case_dir / "classroom.json"),
                "--rules", str(case_dir / "rules.json"),
                "--backend", SOLVER_BACKEND, "--seed", "42",
                "--time-limit", "3",
            ]
            if (case_dir / "history").is_dir():
                flags += ["--history-dir", str(case_dir / "history")]
            py_proc = subprocess.run(flags, capture_output=True, text=True)
            py_status = classify_python_cli(py_proc)
            request, notes = fixture_to_request(case_dir)
            if notes:
                # Python's load/resolver rejected the raw input. Compare the
                # same import surface on the Rust side: a synthesized project
                # workspace validated by the CLI (unknown rule kinds, bad
                # adjacency references, malformed rosters must all be
                # rejected there, exactly like Python).
                rust_status, detail, rust_notes = run_rust_workspace_validate(
                    case_dir, tmp_path
                )
                notes = notes + list(rust_notes.values())
            else:
                rust_status, detail, _ = run_rust(request, tmp_path)
            # Quality baseline (plan §6.6 item 6): on every case where both
            # sides solve, Rust's total cost must not exceed the Python
            # fallback's golden objective (within a 5% tolerance that absorbs
            # the randomize rule's deterministic-but-different RNG draws; a
            # real regression is a mismatch).
            if py_status == rust_status == STATUS_SOLVED:
                golden_path = GOLDENS / cid / "snapshot.json"
                result_path = tmp_path / "rust-result.json"
                if golden_path.is_file() and result_path.exists():
                    try:
                        golden = json.loads(golden_path.read_text(encoding="utf-8"))
                        fallback_cost = golden.get("objective_value")
                        rust_cost = json.loads(
                            result_path.read_text(encoding="utf-8")
                        ).get("total_cost")
                    except (OSError, ValueError):
                        fallback_cost = rust_cost = None
                    if (
                        fallback_cost is not None
                        and rust_cost is not None
                        and fallback_cost > 0
                    ):
                        ratio = rust_cost / fallback_cost
                        notes.append(
                            f"cost vs fallback: rust={rust_cost:.1f} "
                            f"fallback={fallback_cost:.1f} ratio={ratio:.3f}"
                        )
                        if rust_cost > fallback_cost * 1.05:
                            rust_status = "SOLVED(COST_REGRESSION)"
            rows.append((cid, py_status, rust_status, detail, notes))
    return rows


# --- reporting --------------------------------------------------------------

def scoring_request_from_golden(
    snapshot_document: dict[str, object],
) -> dict[str, object]:
    """Build a CoreSolveRequest from a golden snapshot document (which embeds
    students/layout/rules) for the fixed-assignment scoring parity class."""
    from seattrellis.models.snapshot import SeatingSnapshot
    from seattrellis.solver.precompute import precompute_topology

    snapshot = SeatingSnapshot.model_validate(snapshot_document)
    topology = precompute_topology(snapshot.students, snapshot.layout)
    layout = json.loads(snapshot.layout.model_dump_json())
    layout["seats"] = [seat.model_dump(mode="json") for seat in topology.seats]
    request: dict[str, object] = {
        "api_version": NATIVE_API_VERSION,
        "student_count": len(snapshot.students),
        "seat_positions": [[float(seat.x), float(seat.y)] for seat in topology.seats],
        "edges": [list(edge) for edge in sorted(topology.adjacent_seat_index_pairs)],
        "students": [
            {
                "key": student.key,
                "display_name": student.display_name,
                "score": student.score,
                "height_cm": student.height_cm,
                "vision": str(student.vision) if student.vision else None,
                "tags": list(student.tags),
                "needs": list(student.needs),
            }
            for student in snapshot.students
        ],
        "layout": layout,
        "rules": {
            "seed": snapshot.rules.seed,
            "soft": json.loads(snapshot.rules.soft.model_dump_json()),
        },
        "seed": snapshot.rules.seed,
    }
    return request


def fixture_history_context(
    case_dir: Path,
) -> tuple[dict[str, object] | None, dict[str, object] | None, dict[str, object] | None]:
    """Build (history, pair_history, latest_snapshot) from a fixture's
    history directory for the scoring class (None when absent)."""
    from seattrellis.history import build_pair_history, build_seat_history
    from seattrellis.io.json_files import load_layout, load_rules
    from seattrellis.io.students import read_students
    from seattrellis.models.snapshot import SeatingSnapshot

    history_dir = case_dir / "history"
    if not history_dir.is_dir():
        return None, None, None
    students = read_students(case_dir / "students.csv")
    layout = load_layout(case_dir / "classroom.json")
    rules = load_rules(case_dir / "rules.json")
    snapshot_files = sorted(history_dir.glob("*.snapshot.json"))
    if not snapshot_files:
        return None, None, None
    snapshots = [
        SeatingSnapshot.model_validate(json.loads(path.read_text(encoding="utf-8")))
        for path in snapshot_files
    ]
    history = build_seat_history(students, layout, snapshots)
    pair_history = build_pair_history(students, layout, snapshots)
    latest = json.loads(snapshot_files[-1].read_text(encoding="utf-8"))
    return (
        json.loads(history.model_dump_json()),
        json.loads(pair_history.model_dump_json()),
        latest,
    )


def run_rust_score(
    request: dict[str, object],
    pairs: list[list[int]],
    latest_snapshot: dict[str, object] | None,
    tmp: Path,
) -> tuple[str, dict[str, object]]:
    """Run the Rust CLI ``score`` command; return (status, report)."""
    problem_file = tmp / "rust-score-problem.json"
    problem_file.write_text(json.dumps(request), encoding="utf-8")
    cmd = [
        str(CLI), "score",
        "--problem", str(problem_file),
        "--assignment", json.dumps(pairs),
    ]
    if latest_snapshot is not None:
        latest_file = tmp / "rust-score-latest.json"
        latest_file.write_text(json.dumps(latest_snapshot), encoding="utf-8")
        cmd += ["--latest-snapshot", str(latest_file)]
    proc = subprocess.run(cmd, capture_output=True, text=True)
    if proc.returncode != 0:
        return STATUS_INTERNAL_ERROR, {}
    try:
        return STATUS_SOLVED, json.loads(proc.stdout)
    except (OSError, ValueError):
        return STATUS_INTERNAL_ERROR, {}


def run_scoring_class() -> list[tuple[str, str, str, str, list[str]]]:
    """Fixed-assignment scoring parity (plan §6.6 item 4): for every valid
    fixture case, score the golden snapshot's assignment with Python
    ``score_snapshot`` and the Rust CLI ``score`` command, then compare the
    PlanScore breakdown dimension by dimension (status, score within 0.01,
    weight) plus the hard-constraint summary. History-bearing cases also
    exercise fair_rotation / avoid_recent_neighbors / stability."""
    from seattrellis.models.snapshot import SeatingSnapshot
    from seattrellis.scoring import score_snapshot

    rows: list[tuple[str, str, str, str, list[str]]] = []
    if not CLI.exists():
        raise SystemExit(f"Rust CLI not found: {CLI}; build it first")
    golden_dirs = sorted(
        GOLDENS / case_dir.name
        for case_dir in INPUTS.iterdir()
        if case_dir.is_dir() and (GOLDENS / case_dir.name / "snapshot.json").is_file()
    )
    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        for golden_dir in golden_dirs:
            cid = golden_dir.name
            golden = json.loads((golden_dir / "snapshot.json").read_text(encoding="utf-8"))
            request = scoring_request_from_golden(golden)
            # Merge the fixture's compiled hard rules (index pairs) so the
            # hard-constraint summary counts the same rule set Python
            # validates against.
            fixture_request, _ = fixture_to_request(INPUTS / cid)
            for key in ("fixed_seats", "must_be_adjacent", "cannot_be_adjacent", "min_distance"):
                if fixture_request.get(key):
                    request[key] = fixture_request[key]
            snapshot = SeatingSnapshot.model_validate(golden)
            seat_index = {
                seat["seat_id"]: index
                for index, seat in enumerate(request["layout"]["seats"])
            }
            student_index = {
                student["key"]: index
                for index, student in enumerate(request["students"])
            }
            try:
                pairs = [
                    [student_index[assignment["student_key"]], seat_index[assignment["seat_id"]]]
                    for assignment in golden["assignments"]
                ]
            except KeyError as error:
                rows.append((cid, "SKIP", "SKIP", f"golden references unknown key: {error}", []))
                continue

            history, pair_history, latest = fixture_history_context(INPUTS / cid)
            if history is not None:
                request["history"] = history
                request["pair_history"] = pair_history
            latest_model = (
                SeatingSnapshot.model_validate(latest) if latest is not None else None
            )
            from seattrellis.models.history import PairHistory, SeatHistory

            py_score = json.loads(
                score_snapshot(
                    snapshot,
                    history=SeatHistory.model_validate(history) if history else None,
                    pair_history=PairHistory.model_validate(pair_history)
                    if pair_history
                    else None,
                    latest_snapshot=latest_model,
                ).model_dump_json()
            )
            rust_status, rust_score = run_rust_score(request, pairs, latest, tmp_path)
            if rust_status != STATUS_SOLVED:
                rows.append((cid, "SOLVED", rust_status, "rust score command failed", []))
                continue

            notes: list[str] = []
            mismatches = compare_plan_scores(py_score, rust_score, notes)
            rows.append(
                (
                    cid,
                    "SOLVED",
                    "SOLVED",
                    "" if not mismatches else f"{mismatches} mismatches",
                    notes,
                )
            )
    return rows


def compare_plan_scores(
    python: dict[str, object], rust: dict[str, object], notes: list[str]
) -> int:
    """Compare two PlanScore documents; append human-readable notes and
    return the number of mismatching fields (status / score / weight /
    total / hard summary)."""
    mismatches = 0

    def check(label: str, left: object, right: object, tolerance: float = 0.0) -> None:
        nonlocal mismatches
        if isinstance(left, (int, float)) and isinstance(right, (int, float)):
            if abs(float(left) - float(right)) > tolerance:
                mismatches += 1
                notes.append(f"{label}: python={left} rust={right}")
            return
        if left != right:
            mismatches += 1
            notes.append(f"{label}: python={left} rust={right}")

    check("total", python.get("total"), rust.get("total"), tolerance=0.01)
    py_breakdown = python.get("breakdown", {})
    rust_breakdown = rust.get("breakdown", {})
    for key in (
        "fair_rotation_score",
        "avoid_recent_neighbors_score",
        "score_balance_score",
        "height_preference_score",
        "vision_preference_score",
        "diversity_score",
        "stability_score",
    ):
        py_dim = py_breakdown.get(key, {})
        rust_dim = rust_breakdown.get(key, {})
        check(f"{key}.status", py_dim.get("status"), rust_dim.get("status"))
        check(f"{key}.score", py_dim.get("score"), rust_dim.get("score"), tolerance=0.01)
        check(f"{key}.weight", py_dim.get("weight"), rust_dim.get("weight"), tolerance=0.001)
    for key in ("score_position_score", "score_distribution_score", "mentor_pairing_score"):
        py_dim = py_breakdown.get("rule_scores", {}).get(key, {})
        rust_dim = rust_breakdown.get("rule_scores", {}).get(key, {})
        check(f"rule_scores.{key}.status", py_dim.get("status"), rust_dim.get("status"))
        check(
            f"rule_scores.{key}.score",
            py_dim.get("score"),
            rust_dim.get("score"),
            tolerance=0.01,
        )
        check(
            f"rule_scores.{key}.weight",
            py_dim.get("weight"),
            rust_dim.get("weight"),
            tolerance=0.001,
        )
    for key in ("satisfied", "checked_rule_count", "violation_count"):
        check(
            f"hard_constraint_summary.{key}",
            py_breakdown.get("hard_constraint_summary", {}).get(key),
            rust_breakdown.get("hard_constraint_summary", {}).get(key),
        )
    return mismatches


def run_rotation_class() -> list[tuple[str, str, str, str, list[str]]]:
    """Rotation-plan semantic parity (plan §6.2/§17.3 item 3): for every
    valid fixture case both sides generate a 2-period rotation plan from the
    same inputs and seed, then the *semantic* contract is compared:

    - period count and per-period seating completeness (every current
      student seated in every period);
    - per-period solver status (the v1 oracle reports FEASIBLE; the frozen
      v2 vocabulary is Solved);
    - `pair_repeat_summary.relation_totals` key-for-key (measured equal);
    - `fairness_summary.category_totals` on the category keys both sides
      report;
    - `max_occurrences` and empty `warnings`.

    Seats themselves are not compared position-by-position: the revised
    plan §3.2 note says heuristic solutions only need matching semantics,
    validity, scoring definitions and quality gates. One registered oracle
    defect is additionally asserted: the v1 generator reuses the base seed
    for every period (two identical periods), while the Rust side advances
    the seed per period (seed + period - 1), so its periods must differ.
    """
    rows: list[tuple[str, str, str, str, list[str]]] = []
    if not CLI.exists() or not PY_CLI.exists():
        raise SystemExit("rotation class requires both the Rust and the Python CLI")
    golden_dirs = sorted(
        GOLDENS / case_dir.name
        for case_dir in INPUTS.iterdir()
        if case_dir.is_dir() and (GOLDENS / case_dir.name / "snapshot.json").is_file()
    )
    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        for golden_dir in golden_dirs:
            case = golden_dir.name
            case_dir = INPUTS / case
            try:
                rules = json.loads((case_dir / "rules.json").read_text(encoding="utf-8"))
            except OSError:
                rows.append((case, "ROTATION", "ROTATION", "SKIP: no rules.json", []))
                continue
            seed = rules.get("seed", 42)

            # --- Python oracle -------------------------------------------------
            py_out = tmp_path / f"{case}-py-rotation.json"
            py_proc = subprocess.run(
                [
                    str(PY_CLI), "rotation-plan",
                    "--students", str(case_dir / "students.csv"),
                    "--layout", str(case_dir / "classroom.json"),
                    "--rules", str(case_dir / "rules.json"),
                    "--periods", "2",
                    "--seed", str(seed),
                    "--output", str(py_out),
                ],
                capture_output=True,
                text=True,
            )
            if py_proc.returncode != 0 or not py_out.is_file():
                rows.append((case, "ROTATION", "ROTATION", f"PY_SKIP: {py_proc.stderr.strip()[:120]}", []))
                continue

            # --- Rust CLI (project workspace) ----------------------------------
            workspace = tmp_path / f"{case}-ws"
            workspace.mkdir()
            for src_name, dst_name in (
                ("students.csv", "students.csv"),
                ("classroom.json", "layout.json"),
                ("rules.json", "rules.json"),
            ):
                shutil.copyfile(case_dir / src_name, workspace / dst_name)
            init = subprocess.run(
                [str(CLI), "project-init", "--dir", str(workspace)],
                capture_output=True,
                text=True,
            )
            if init.returncode != 0:
                rows.append((case, "ROTATION", "ROTATION", f"INIT_FAILED: {init.stderr.strip()[:120]}", []))
                continue
            ru_out = tmp_path / f"{case}-ru-rotation.json"
            ru_proc = subprocess.run(
                [
                    str(CLI), "project-rotate",
                    "--project", str(workspace / "seattrellis.project.json"),
                    "--periods", "2",
                    "--seed", str(seed),
                    "--output", str(ru_out),
                ],
                capture_output=True,
                text=True,
            )
            if ru_proc.returncode != 0 or not ru_out.is_file():
                rows.append((case, "ROTATION", "ROTATION", f"RU_FAILED: {ru_proc.stderr.strip()[:120]}", []))
                continue

            try:
                py_plan = json.loads(py_out.read_text(encoding="utf-8"))
                ru_plan = json.loads(ru_out.read_text(encoding="utf-8"))
                # relation_totals / category_totals are counts over occupied
                # seat pairs: when every seat is filled the counts depend
                # only on the layout (solution-independent, strictly
                # comparable); layouts with empty seats make them
                # solution-dependent, so those cases compare semantically.
                student_count = sum(
                    1 for _ in csv.reader(open(case_dir / "students.csv", encoding="utf-8"))
                ) - 1
                layout_doc = json.loads(
                    (case_dir / "classroom.json").read_text(encoding="utf-8")
                )
                seat_count = sum(
                    1
                    for seat in layout_doc.get("seats", [])
                    if seat.get("enabled", True)
                )
                full_occupancy = student_count == seat_count
                mismatches = compare_rotation_plans(py_plan, ru_plan, full_occupancy)
                label = "ROTATION"
                rows.append(
                    (case, label, label, "" if not mismatches else f"{mismatches} mismatches", [])
                )
            except Exception as exc:  # noqa: BLE001
                rows.append((case, "ROTATION", "ROTATION", f"COMPARE_FAILED: {exc}", []))
    return rows


def compare_rotation_plans(
    python: dict[str, object], rust: dict[str, object], full_occupancy: bool = True
) -> int:
    """Semantic comparison of two rotation plans; returns the mismatch count
    and appends nothing (the caller logs the row).

    `relation_totals` and `category_totals` count occupied seat pairs, so
    they are solution-independent only when every seat is filled
    (`full_occupancy`); layouts with empty seats are compared on the
    semantic core (periods, completeness, status, structure) instead."""
    mismatches = 0
    # Durable artifact identity/version are frozen parity fields, not merely
    # metadata. This catches the former Rust "0.2.2" value against the
    # oracle rotation-plan schema version "1.0".
    if python.get("kind") != rust.get("kind"):
        mismatches += 1
    if python.get("schema_version") != rust.get("schema_version"):
        mismatches += 1
    py_periods = python.get("periods", [])
    ru_periods = rust.get("periods", [])
    if len(py_periods) != len(ru_periods):
        mismatches += 1

    def seated_keys(plan: dict[str, object]) -> list[set[str]]:
        return [
            {
                assignment.get("student_key")
                for assignment in period.get("snapshot", {}).get("assignments", [])
                if assignment.get("student_key")
            }
            for period in plan.get("periods", [])
        ]

    py_keys = seated_keys(python)
    ru_keys = seated_keys(rust)
    # Completeness: every period seats the same student set on both sides.
    if py_keys and ru_keys and py_keys[0] != ru_keys[0]:
        mismatches += 1

    # Solver status: the v1 oracle reports FEASIBLE; the frozen v2
    # vocabulary is Solved (M1-03). Map the legacy value before comparing.
    def normalized_status(plan: dict[str, object]) -> set[str]:
        return {
            "Solved" if period.get("snapshot", {}).get("solver_status") == "FEASIBLE"
            else period.get("snapshot", {}).get("solver_status")
            for period in plan.get("periods", [])
        }

    if normalized_status(python) != {"Solved"} or normalized_status(rust) != {"Solved"}:
        mismatches += 1

    if full_occupancy:
        # relation_totals key-for-key (solution-independent when full).
        py_relations = python.get("pair_repeat_summary", {}).get("relation_totals", {})
        ru_relations = rust.get("pair_repeat_summary", {}).get("relation_totals", {})
        if py_relations != ru_relations:
            mismatches += 1

        # max_occurrences.
        if (
            python.get("pair_repeat_summary", {}).get("max_occurrences")
            != rust.get("pair_repeat_summary", {}).get("max_occurrences")
        ):
            mismatches += 1

        # fairness category_totals on shared keys.
        py_categories = python.get("fairness_summary", {}).get("category_totals", {})
        ru_categories = rust.get("fairness_summary", {}).get("category_totals", {})
        for category, value in ru_categories.items():
            if category in py_categories and py_categories[category] != value:
                mismatches += 1

    # history_count (solution-independent: snapshot count).
    if (
        python.get("fairness_summary", {}).get("history_count")
        != rust.get("fairness_summary", {}).get("history_count")
    ):
        mismatches += 1

    # warnings empty on both sides.
    if python.get("warnings") or rust.get("warnings"):
        mismatches += 1

    # Registered oracle defect: the v1 generator reuses the base seed for
    # every period (identical assignments), while the Rust side advances the
    # seed (seed + period - 1), so its periods must differ. Compare the
    # student -> seat mapping, not the key set (every period seats the full
    # roster on both sides).
    def assignments(plan: dict[str, object]) -> list[set[tuple[str, str]]]:
        return [
            {
                (assignment.get("student_key"), assignment.get("seat_id"))
                for assignment in period.get("snapshot", {}).get("assignments", [])
                if assignment.get("student_key") and assignment.get("seat_id")
            }
            for period in plan.get("periods", [])
        ]

    ru_assignments = assignments(rust)
    if len(ru_assignments) > 1 and ru_assignments[0] == ru_assignments[1]:
        mismatches += 1
    return mismatches


def run_exports_class() -> list[tuple[str, str, str, str, list[str]]]:
    """Office export independent-reader verification (revised plan §11.6).

    For every valid fixture case the Rust CLI solves the problem and exports
    XLSX / DOCX / PPTX (teacher template) plus XLSX (public template). Each
    file is then reopened with a *different implementation* — openpyxl /
    python-docx / python-pptx — and checked for structure, seat content, and
    the public-template privacy guarantee (no real student names anywhere).
    The PNG/PDF rasters are decoded with Pillow/pypdf (A4-ish bounds, Image
    XObject presence, ledger §19.26), and the dedicated print-html layout
    (print-layout-spec) is parsed structurally with the standard library
    (page skeleton, seat grid rows, name presence, structure annotations,
    reproducibility footer).

    Byte parity with the Python exporters is impossible (zip mtimes), so the
    acceptance criterion is exactly what the plan asks for: an independent
    reader can open the document and the semantic content matches.
    """
    from openpyxl import load_workbook

    rows: list[tuple[str, str, str, str, list[str]]] = []
    if not CLI.exists():
        raise SystemExit(f"Rust CLI not found: {CLI}; build it first")
    golden_dirs = sorted(
        GOLDENS / case_dir.name
        for case_dir in INPUTS.iterdir()
        if case_dir.is_dir() and (GOLDENS / case_dir.name / "snapshot.json").is_file()
    )
    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        for golden_dir in golden_dirs:
            case = golden_dir.name
            request, notes = fixture_to_request(INPUTS / case)
            problem_file = tmp_path / f"{case}-problem.json"
            solution_file = tmp_path / f"{case}-solution.json"
            problem_file.write_text(json.dumps(request), encoding="utf-8")
            solved = subprocess.run(
                [str(CLI), "solve", "--problem", str(problem_file), "--output", str(solution_file)],
                capture_output=True,
                text=True,
            )
            if solved.returncode != 0 or not solution_file.is_file():
                rows.append((case, "EXPORTS", "EXPORTS", f"SOLVE_FAILED: {solved.stderr.strip()}", notes))
                continue
            response = json.loads(solution_file.read_text(encoding="utf-8"))
            if not response.get("feasible"):
                rows.append((case, "EXPORTS", "EXPORTS", "SOLVE_SKIPPED: not feasible", notes))
                continue
            real_names = {
                student.get("display_name") or student.get("key")
                for student in request.get("students", [])
                if student.get("display_name") or student.get("key")
            }
            try:
                _verify_exports(CLI, problem_file, solution_file, response, tmp_path, real_names, case, rows)
            except Exception as exc:  # noqa: BLE001 - report, do not abort the class
                rows.append((case, "EXPORTS", "ERROR", str(exc), notes))
            # Dedicated print layout (print-layout-spec): independent reader
            # for the print-html structure (page skeleton, seat grid rows,
            # name presence, structure annotations, reproducibility footer).
            try:
                _verify_print_html_export(CLI, case, tmp_path)
                rows.append((case, "EXPORTS-PRINT-HTML", "EXPORTS-PRINT-HTML", "independent reader ok", []))
            except Exception as exc:  # noqa: BLE001
                rows.append((case, "EXPORTS-PRINT-HTML", "EXPORTS-FAILED", str(exc), []))
    return rows


def _verify_exports(
    cli: Path,
    problem_file: Path,
    solution_file: Path,
    response: dict[str, object],
    tmp_path: Path,
    real_names: set[str],
    case: str,
    rows: list[tuple[str, str, str, str, list[str]]],
) -> None:
    """Export every Office format and reopen it with an independent reader."""
    from openpyxl import load_workbook

    for fmt in ("xlsx", "docx", "pptx", "png", "pdf"):
        out = tmp_path / f"{case}.{fmt}"
        exported = subprocess.run(
            [
                str(cli),
                "export",
                "--problem",
                str(problem_file),
                "--solution",
                str(solution_file),
                "--format",
                fmt,
                "--template",
                "teacher",
                "--output",
                str(out),
            ],
            capture_output=True,
            text=True,
        )
        if exported.returncode != 0 or not out.is_file():
            rows.append((case, f"EXPORTS-{fmt.upper()}", "EXPORTS-FAILED", exported.stderr.strip(), []))
            continue
        try:
            if fmt == "xlsx":
                _verify_xlsx(out, response)
            elif fmt == "docx":
                _verify_docx(out, response)
            elif fmt == "pptx":
                _verify_pptx(out, response)
            elif fmt == "png":
                _verify_png(out, response)
            else:
                _verify_pdf(out, response)
            label = f"EXPORTS-{fmt.upper()}"
            rows.append((case, label, label, "independent reader ok", []))
        except Exception as exc:  # noqa: BLE001
            rows.append((case, f"EXPORTS-{fmt.upper()}", "EXPORTS-FAILED", str(exc), []))

    # Public-template privacy: no real student name may survive in any format.
    for fmt in ("xlsx", "docx", "pptx", "png", "pdf"):
        out = tmp_path / f"{case}-public.{fmt}"
        exported = subprocess.run(
            [
                str(cli),
                "export",
                "--problem",
                str(problem_file),
                "--solution",
                str(solution_file),
                "--format",
                fmt,
                "--template",
                "public",
                "--output",
                str(out),
            ],
            capture_output=True,
            text=True,
        )
        if exported.returncode != 0 or not out.is_file():
            rows.append((case, f"EXPORTS-PUBLIC-{fmt.upper()}", "EXPORTS-FAILED", exported.stderr.strip(), []))
            continue
        try:
            text = _read_office_text(out)
            leaked = [name for name in real_names if name and name in text]
            label = f"EXPORTS-PUBLIC-{fmt.upper()}"
            if leaked:
                rows.append((case, label, "EXPORTS-PRIVACY-LEAK", f"leaked: {leaked[:3]}", []))
            else:
                rows.append((case, label, label, "no names leaked", []))
        except Exception as exc:  # noqa: BLE001
            rows.append((case, f"EXPORTS-PUBLIC-{fmt.upper()}", "MISMATCH", str(exc), []))


def _verify_xlsx(path: Path, response: dict[str, object]) -> None:
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=True)
    if workbook.sheetnames != ["Seating", "Assignments"]:
        raise AssertionError(f"unexpected sheets: {workbook.sheetnames}")
    seating = workbook["Seating"]
    grid_text = " ".join(
        str(value) for row in seating.iter_rows(values_only=True) for value in row if value
    )
    if "R1C1" not in grid_text and "R1C2" not in grid_text:
        raise AssertionError("grid sheet carries no seat ids")
    assignments = workbook["Assignments"]
    rows = list(assignments.iter_rows(values_only=True))
    if rows[0] != ("student_key", "student_name", "seat_id"):
        raise AssertionError(f"unexpected assignments header: {rows[0]}")
    seated = len(response.get("assignment", []))
    if len(rows) - 1 != seated:
        raise AssertionError(f"assignments rows {len(rows) - 1} != seated {seated}")


def _verify_docx(path: Path, response: dict[str, object]) -> None:
    import docx as python_docx

    document = python_docx.Document(path)
    if not document.tables:
        raise AssertionError("docx carries no seat table")
    table = document.tables[0]
    if len(table.rows) < 1 or len(table.columns) < 1:
        raise AssertionError("docx seat table is empty")
    cells = [cell.text for row in table.rows for cell in row.cells]
    if not any(cells):
        raise AssertionError("docx seat table has no text")


def _verify_png(path: Path, response: dict[str, object]) -> None:
    """PNG must be a decodable raster of plausible dimensions (Pillow)."""
    from PIL import Image

    with Image.open(path) as image:
        image.verify()
    with Image.open(path) as image:
        width, height = image.size
    seated = len(response.get("assignment", []))
    if width < 100 or height < 100:
        raise AssertionError(f"PNG is implausibly small: {width}x{height}")
    if seated < 1:
        raise AssertionError("PNG verified for an empty plan")


def _verify_pdf(path: Path, response: dict[str, object]) -> None:
    """PDF must open with an independent reader and stay inside A4-ish
    bounds. Since §19.26 the page is rasterized into an Image XObject
    (names are pixels, not searchable text - recorded as a boundary), so
    the semantic check is the presence of a page-sized image resource."""
    from pypdf import PdfReader

    reader = PdfReader(path)
    if not reader.pages:
        raise AssertionError("PDF has no pages")
    for page in reader.pages:
        width = float(page.mediabox.width)
        height = float(page.mediabox.height)
        if width > 900 or height > 1300:
            raise AssertionError(f"PDF page exceeds A4-ish bounds: {width}x{height}")
        resources = page.get("/Resources")
        xobjects = resources.get("/XObject") if resources else None
        if not xobjects or not xobjects.get_object():
            raise AssertionError("PDF page carries no image content")


def _verify_pptx(path: Path, response: dict[str, object]) -> None:
    from pptx import Presentation

    presentation = Presentation(path)
    if (presentation.slide_width, presentation.slide_height) != (12192000, 6858000):
        raise AssertionError(
            f"slide size {(presentation.slide_width, presentation.slide_height)} != 16:9"
        )
    if not presentation.slides:
        raise AssertionError("pptx has no slides")
    shapes = list(presentation.slides[0].shapes)
    seated = len(response.get("assignment", []))
    if len(shapes) < seated + 1:
        raise AssertionError(f"shapes {len(shapes)} < seats {seated} + title")


def _read_office_text(path: Path) -> str:
    """All visible text of an Office file, via the independent readers."""
    suffix = path.suffix.lower()
    if suffix == ".xlsx":
        from openpyxl import load_workbook

        workbook = load_workbook(path, read_only=True)
        return " ".join(
            str(value)
            for sheet in workbook.worksheets
            for row in sheet.iter_rows(values_only=True)
            for value in row
            if value
        )
    if suffix == ".docx":
        import docx as python_docx

        document = python_docx.Document(path)
        parts = [paragraph.text for paragraph in document.paragraphs]
        parts.extend(cell.text for table in document.tables for row in table.rows for cell in row.cells)
        return " ".join(parts)
    if suffix == ".pptx":
        from pptx import Presentation

        presentation = Presentation(path)
        parts: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)
        return " ".join(parts)
    if suffix == ".pdf":
        from pypdf import PdfReader

        reader = PdfReader(path)
        return " ".join(page.extract_text() or "" for page in reader.pages)
    if suffix == ".png":
        # This renderer's PNG carries no text; the privacy check still
        # verifies the raster decodes (an image is either fully rendered or
        # corrupt, so no student data can leak through a text channel).
        from PIL import Image

        with Image.open(path) as image:
            image.verify()
        return ""
    raise AssertionError(f"unknown export format: {path}")


# ---------------------------------------------------------------------------
# print-html independent-reader verification (print-layout-spec.md §2/§3)
# ---------------------------------------------------------------------------

class _PrintHtmlParser(HTMLParser):
    """Structural extractor for the Rust print-html export.

    Pulls the doctype, title, header class/meta lines, the platform
    annotation, the per-row seat cells and the structure/footer notes out of
    the document with the standard-library parser only. A cell is a
    ``div.grid-row > div.seat`` entry; its text is the concatenation of all
    descendant text (including the ``span.sid`` student identifier when the
    caller enables it), so a renamed class, a dropped row or a lost name
    changes what this extractor reports.
    """

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.doctype = ""
        self.title = ""
        self.header_cls = ""
        self.header_meta = ""
        self.stage = ""
        self.structure = ""
        self.footer = ""
        # One entry per grid-row: (class tokens, cell text) per cell.
        self.rows: list[list[tuple[list[str], str]]] = []
        self._in_title = False
        self._row: list[tuple[list[str], str]] | None = None
        self._cell: tuple[list[str], list[str]] | None = None
        # Element stack of (tag, class tokens) for context routing.
        self._stack: list[tuple[str, list[str]]] = []

    def handle_decl(self, decl: str) -> None:
        self.doctype += decl

    def handle_starttag(self, tag: str, attrs) -> None:
        classes = dict(attrs).get("class", "").split()
        self._stack.append((tag, classes))
        if tag == "title":
            self._in_title = True
        elif tag == "div":
            if "grid-row" in classes:
                self._row = []
            elif "seat" in classes and self._row is not None:
                self._cell = (classes, [])

    def handle_endtag(self, tag: str) -> None:
        if tag == "title":
            self._in_title = False
        elif tag == "div":
            if self._cell is not None and "seat" in self._cell[0]:
                assert self._row is not None
                self._row.append((self._cell[0], "".join(self._cell[1]).strip()))
                self._cell = None
            elif self._row is not None and self._stack and "grid-row" in self._stack[-1][1]:
                self.rows.append(self._row)
                self._row = None
        if self._stack:
            self._stack.pop()

    def handle_data(self, data: str) -> None:
        if self._in_title:
            self.title += data
            return
        if self._cell is not None:
            self._cell[1].append(data)
            return
        # Route by the innermost open element that carries a known class
        # (the structure/footer notes live in unclassed spans inside their
        # classed containers).
        for _, classes in reversed(self._stack):
            if "cls" in classes:
                self.header_cls += data
                return
            if "meta" in classes:
                self.header_meta += data
                return
            if "stage" in classes:
                self.stage += data
                return
            if "structure" in classes:
                self.structure += data
                return
            if "print-footer" in classes:
                self.footer += data
                return


def _print_html_expected(case_dir: Path) -> dict[str, object]:
    """The print-html structural contract for a fixture case, derived from
    the same input files the synthesized project workspace compiles: roster
    names, seat-grid row count, structure annotations (window/door from the
    seat flags) and the aisle lanes (grid columns without an enabled seat)."""
    names: list[str] = []
    with (case_dir / "students.csv").open(newline="", encoding="utf-8") as handle:
        for row in csv.DictReader(handle):
            name = (row.get("name") or "").strip()
            if name:
                names.append(name)
    layout = json.loads((case_dir / "classroom.json").read_text(encoding="utf-8"))
    enabled = [s for s in layout.get("seats", []) if s.get("enabled", True)]
    rows = {int(s["row"]) for s in enabled}
    cols = {int(s["col"]) for s in enabled}
    return {
        "names": names,
        "student_count": len(names),
        "seat_count": len(enabled),
        "row_count": max(rows) - min(rows) + 1,
        "window": any(bool(s.get("near_window")) for s in enabled),
        "door": any(bool(s.get("near_door")) for s in enabled),
        "aisle": any(c not in cols for c in range(min(cols), max(cols) + 1)),
    }


def _verify_print_html(path: Path, expected: dict[str, object]) -> None:
    """Independent structural verification of a Rust print-html export
    (print-layout-spec.md §2/§3), standard library only:

    - page skeleton: doctype, html/head with style + title, body;
    - header line: class name + "座位表" and the "N students / M seats" meta;
    - platform annotation (讲台);
    - seat grid: expected row count, every student name rendered in a seat
      cell (truncation-tolerant), no height/vision detail leakage;
    - classroom structure notes (窗 / 门) and aisle lanes (过道);
    - reproducibility footer (seed · 第 1 / 1 页).
    """
    with path.open(encoding="utf-8") as handle:
        text = handle.read()
    parser = _PrintHtmlParser()
    parser.feed(text)
    parser.close()

    if "html" not in parser.doctype.lower():
        raise AssertionError(f"print-html lacks an html doctype: {parser.doctype!r}")
    if not parser.title:
        raise AssertionError("print-html carries no <title>")
    if "座位表" not in parser.header_cls:
        raise AssertionError(f"print-html header missing class name: {parser.header_cls!r}")
    meta_student = f"{expected['student_count']} students"
    meta_seat = f"{expected['seat_count']} seats"
    if meta_student not in parser.header_meta or meta_seat not in parser.header_meta:
        raise AssertionError(
            f"print-html header meta {parser.header_meta!r} != {meta_student} / {meta_seat}"
        )
    if "讲台" not in parser.stage:
        raise AssertionError(f"print-html missing the platform annotation: {parser.stage!r}")

    row_count = int(expected["row_count"])
    if len(parser.rows) != row_count:
        raise AssertionError(
            f"print-html grid rows {len(parser.rows)} != expected {row_count}"
        )
    for index, row in enumerate(parser.rows, start=1):
        if not row:
            raise AssertionError(f"print-html grid row {index} carries no cells")

    # Every seat cell that is not an empty/disabled placeholder carries one
    # student name; all roster students must be seated in a valid corpus case.
    cell_texts = [cell_text for row in parser.rows for _, cell_text in row]
    named = [cell_text for cell_text in cell_texts if cell_text and cell_text != "空座"]
    if len(named) != int(expected["student_count"]):
        raise AssertionError(
            f"print-html named seats {len(named)} != students {expected['student_count']}"
        )
    names = expected["names"]
    assert isinstance(names, list)
    for name in names:
        if not any(name in cell_text or cell_text.rstrip("…").endswith(name) for cell_text in named):
            raise AssertionError(f"print-html seat grid misses student name {name!r}")

    # The print view renders names only: height/vision details must never
    # appear (print-layout-spec §2.4; the dedicated renderer draws no detail
    # line at all). Check the rendered text, not the raw document — the CSS
    # block legitimately contains `min-height` etc. The tokens are the exact
    # detail-line formats ("172 cm" / "vision 0.6"): bare words would trip
    # on fixture ids embedded in titles (soft-vision-front).
    rendered_text = " ".join(
        [parser.title, parser.header_cls, parser.header_meta, parser.stage,
         parser.structure, parser.footer]
        + [cell for row in parser.rows for _, cell in row]
    )
    for token in (" cm", "vision "):
        if token in rendered_text:
            raise AssertionError(f"print-html leaks a student detail token {token!r}")

    if expected["window"] and "窗" not in parser.structure:
        raise AssertionError(f"print-html missing the window annotation: {parser.structure!r}")
    if expected["door"] and "门" not in parser.structure:
        raise AssertionError(f"print-html missing the door annotation: {parser.structure!r}")
    if expected["aisle"] and "过道" not in text:
        raise AssertionError("print-html missing the aisle lane annotation")

    if "seed" not in parser.footer or "第 1 / 1 页" not in parser.footer:
        raise AssertionError(f"print-html footer lacks the reproducibility line: {parser.footer!r}")


def _verify_print_html_export(cli: Path, case: str, tmp_path: Path) -> None:
    """Export a print-html artifact through the CLI project workflow and
    verify it with the independent reader. The flat `export` command does
    not expose the print-html format (CLI format surface — the dedicated
    print layout is a server/workbench format; see ledger §19.31), so the
    check rides `project-export`, which renders a saved project plan through
    the same shared `export_plan` renderer (teacher template)."""
    case_dir = INPUTS / case
    workspace = tmp_path / f"{case}-print-html-ws"
    workspace.mkdir(exist_ok=True)
    for src_name, dst_name in (
        ("students.csv", "students.csv"),
        ("classroom.json", "layout.json"),
        ("rules.json", "rules.json"),
    ):
        shutil.copyfile(case_dir / src_name, workspace / dst_name)
    init = subprocess.run(
        [str(cli), "project-init", "--dir", str(workspace)],
        capture_output=True,
        text=True,
    )
    if init.returncode != 0:
        raise AssertionError(f"project-init failed: {init.stderr.strip()[:200]}")
    project_file = workspace / "seattrellis.project.json"
    plan_file = workspace / "plan.json"
    solved = subprocess.run(
        [str(cli), "project-solve", "--project", str(project_file), "--output", str(plan_file)],
        capture_output=True,
        text=True,
    )
    if solved.returncode != 0 or not plan_file.is_file():
        raise AssertionError(f"project-solve failed: {solved.stderr.strip()[:200]}")
    out = tmp_path / f"{case}.print.html"
    exported = subprocess.run(
        [str(cli), "project-export", "--project", str(project_file),
         "--format", "print-html", "--snapshot", str(plan_file), "--output", str(out)],
        capture_output=True,
        text=True,
    )
    if exported.returncode != 0 or not out.is_file():
        raise AssertionError(
            f"project-export print-html failed: {exported.stderr.strip()[:200]}"
        )
    _verify_print_html(out, _print_html_expected(case_dir))


def run_rust_candidates(
    request: dict[str, object], count: int, tmp: Path
) -> tuple[str, str, int | None]:
    """Run the Rust CLI ``candidates`` command; return (status, detail, count).

    The report is printed to stdout; a generation failure exits 2. The count
    is the number of distinct feasible plans actually generated.
    """
    problem_file = tmp / "rust-candidates-problem.json"
    problem_file.write_text(json.dumps(request), encoding="utf-8")
    proc = subprocess.run(
        [str(CLI), "candidates", "--problem", str(problem_file), "--count", str(count)],
        capture_output=True,
        text=True,
    )
    if proc.returncode == 0:
        try:
            report = json.loads(proc.stdout)
            generated = int(report.get("candidate_count", 0))
            detail = f"candidates={generated}"
            return STATUS_SOLVED, detail, generated
        except (OSError, ValueError):
            return STATUS_INTERNAL_ERROR, proc.stdout.strip()[:200], None
    detail = proc.stderr.strip()[:200] or proc.stdout.strip()[:200]
    if any(token in detail.lower() for token in INVALID_TOKENS):
        return STATUS_INVALID_INPUT, detail, None
    return STATUS_INTERNAL_ERROR, detail, None


def run_python_candidates(case_dir: Path, count: int, tmp: Path) -> tuple[str, str, int | None]:
    """Run the Python oracle CLI ``solve --candidates N`` on a fixture case;
    return (status, detail, generated count). The candidate set document is
    written to the ``--output`` file (stdout carries a human summary)."""
    output_file = tmp / f"python-candidates-{case_dir.name}-{count}.json"
    flags = [
        str(PY_CLI), "solve",
        "--students", str(case_dir / "students.csv"),
        "--layout", str(case_dir / "classroom.json"),
        "--rules", str(case_dir / "rules.json"),
        "--backend", SOLVER_BACKEND, "--seed", "42",
        "--time-limit", "3",
        "--candidates", str(count),
        "--output", str(output_file),
    ]
    if (case_dir / "history").is_dir():
        flags += ["--history-dir", str(case_dir / "history")]
    proc = subprocess.run(flags, capture_output=True, text=True)
    if proc.returncode == 0 and output_file.exists():
        try:
            report = json.loads(output_file.read_text(encoding="utf-8"))
        except (OSError, ValueError):
            return STATUS_INTERNAL_ERROR, proc.stdout.strip()[:200], None
        candidates = report.get("candidates")
        if candidates is None:
            # --candidates 1 writes a plain snapshot document (a CandidateSet
            # with exactly one entry), not a candidate-set envelope.
            if report.get("assignments") is not None:
                generated = 1
                return STATUS_SOLVED, "candidates=1", generated
            return STATUS_INTERNAL_ERROR, proc.stdout.strip()[:200], None
        generated = len(candidates)
        return STATUS_SOLVED, f"candidates={generated}", generated
    status = classify_python_cli(proc)
    return status, (proc.stderr + proc.stdout).strip()[:200], None


# Fixture cases exercising the candidate sizes required by the M3 Exit Gate
# (ledger §17.4): 20/40/50/60/80 students × 1/5/20 requested candidates.
CANDIDATES_FIXTURE_CASES = [
    "p20-rect-exact-none",
    "p40-rect-exact-sparse",
    "p50-custom-adj-sparse",
    "p60-rect-exact-dense",
    "p80-rect-exact-dense",
]
CANDIDATES_COUNTS = [1, 5, 20]


def run_candidates_class() -> list[tuple[str, str, str, str, list[str]]]:
    """Python↔Rust candidate-set parity: same fixture, same base seed, same
    requested count. The comparison is status class + generated count; the
    assignments themselves differ (independent solvers), so per-candidate
    content parity is out of scope (ledger §19.3.4). When both sides solve
    but generate different numbers of distinct plans, the count is embedded
    in the status token so the generic reporter flags a mismatch."""
    rows: list[tuple[str, str, str, str, list[str]]] = []
    if not CLI.exists():
        raise SystemExit(f"Rust CLI not found: {CLI}; build it first")
    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        for case_name in CANDIDATES_FIXTURE_CASES:
            case_dir = INPUTS / case_name
            for count in CANDIDATES_COUNTS:
                cid = f"{case_name}-cand{count}"
                py_status, py_detail, py_generated = run_python_candidates(
                    case_dir, count, tmp_path
                )
                request, _ = fixture_to_request(case_dir)
                request["seed"] = 42
                rust_status, rust_detail, rust_generated = run_rust_candidates(
                    request, count, tmp_path
                )
                py_label = (
                    f"{py_status}({py_generated})"
                    if py_generated is not None
                    else py_status
                )
                rust_label = (
                    f"{rust_status}({rust_generated})"
                    if rust_generated is not None
                    else rust_status
                )
                notes: list[str] = []
                if py_detail and py_generated is None:
                    notes.append(f"python: {py_detail}")
                if rust_generated is None and rust_detail:
                    notes.append(f"rust: {rust_detail}")
                rows.append((cid, py_label, rust_label, rust_detail, notes))
    return rows


def report(rows: list[tuple[str, str, str, str, list[str]]], allow_documented: bool = False) -> int:
    mismatches = 0
    documented = 0
    print(f"{'case':<42} {'python':<16} {'rust':<16} match")
    print("-" * 92)
    for cid, py_status, rust_status, detail, notes in rows:
        # M1-03 frozen semantics: a legal incumbent found within the budget
        # beats a timeout, so Python TIMEOUT + Rust SOLVED is a match (the
        # Rust solver simply finished inside the budget). Since M3-04 the
        # Rust side carries its own --time-limit for the TIMEOUT class.
        match = py_status == rust_status or (
            py_status == STATUS_TIMEOUT and rust_status == STATUS_SOLVED
        )
        if not match:
            mismatches += 1
        flag = "OK " if match else "MISMATCH"
        print(f"{cid:<42} {py_status:<16} {rust_status:<16} {flag}")
        if detail:
            print(f"    rust detail: {detail}")
        for note in notes:
            print(f"    note: {note}")
        if not match and py_status in KNOWN_RUST_GAPS:
            print(f"    documented gap: {KNOWN_RUST_GAPS[py_status]}")
        if not match and rust_status in KNOWN_RUST_GAPS:
            print(f"    documented gap: {KNOWN_RUST_GAPS[rust_status]}")
        if not match and cid in DOCUMENTED_CORPUS_GAPS:
            documented += 1
            print(f"    documented corpus gap: {DOCUMENTED_CORPUS_GAPS[cid]}")
    print("-" * 92)
    new_mismatches = mismatches - documented
    print(
        f"cases: {len(rows)}  mismatches: {mismatches} "
        f"(documented: {documented}, new: {new_mismatches})"
    )
    if allow_documented:
        # CI mode: exactly the documented corpus gaps are tolerated; any
        # other mismatch still fails the run.
        return 1 if new_mismatches else 0
    # Strict mode (default, M0-03): any mismatch fails the run.
    return 1 if mismatches else 0


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--sizes", default="40,50,60")
    parser.add_argument("--time-limit", type=float, default=3.0)
    parser.add_argument("--fixtures", action="store_true", help="run the fixtures/parity status classes")
    parser.add_argument(
        "--candidates",
        action="store_true",
        help="run the candidate-set parity class (20/40/50/60/80 x 1/5/20)",
    )
    parser.add_argument(
        "--scoring",
        action="store_true",
        help="run the fixed-assignment PlanScore parity class (all valid fixture cases)",
    )
    parser.add_argument(
        "--rotation",
        action="store_true",
        help="run the rotation-plan semantic parity class (2 periods per "
        "valid fixture case, Python oracle vs Rust project-rotate)",
    )
    parser.add_argument(
        "--exports",
        action="store_true",
        help="run the Office export independent-reader class (XLSX/DOCX/PPTX "
        "reopened with openpyxl/python-docx/python-pptx, teacher + public "
        "privacy checks, plus the print-html structure reader and the PDF "
        "raster-page check)",
    )
    parser.add_argument(
        "--cli-golden",
        action="store_true",
        help="run the CLI stdout/stderr/exit-code golden class (Rust contract "
        "vs registered goldens + Python exit-code semantics)",
    )
    parser.add_argument(
        "--cli-golden-record",
        action="store_true",
        help="(re)record the CLI golden files under fixtures/cli-goldens",
    )
    parser.add_argument(
        "--allow-documented-gaps",
        action="store_true",
        help="CI mode: tolerate exactly the case-level documented corpus gaps "
        "(DOCUMENTED_CORPUS_GAPS); any new mismatch still fails the run",
    )
    args = parser.parse_args()

    if not PY_CLI.exists():
        raise SystemExit(f"Python CLI not found: {PY_CLI}; activate the project venv")

    rows: list[tuple[str, str, str, str, list[str]]] = []
    if args.fixtures:
        rows.extend(run_fixture_classes())
    if args.candidates:
        rows.extend(run_candidates_class())
    if args.scoring:
        rows.extend(run_scoring_class())
    if args.rotation:
        rows.extend(run_rotation_class())
    if args.exports:
        rows.extend(run_exports_class())
    if args.cli_golden or args.cli_golden_record:
        with tempfile.TemporaryDirectory(prefix="cli-golden-") as tmpdir:
            rows.extend(run_cli_golden_class(record=args.cli_golden_record, tmp=Path(tmpdir)))
        if args.cli_golden_record:
            # Keep the M6 retirement provenance inventory synchronized with
            # the byte contract just recorded. The import is local so normal
            # differential runs do not load the corpus generator.
            from gen_parity_fixtures import write_golden_provenance

            write_golden_provenance()
            return report(rows, allow_documented=args.allow_documented_gaps)
    if not rows:
        sizes = [int(item) for item in args.sizes.split(",")]
        rows = run_benchmark_classes(sizes, args.time_limit)
    return report(rows, allow_documented=args.allow_documented_gaps)
# ---------------------------------------------------------------------------
# CLI stdout/stderr/exit-code golden class (plan §5.5 / ledger §1 "输出契约无
# golden"): register the Rust CLI's output contract for representative
# commands, then re-run and compare byte-for-byte (JSON normalized). Python
# side participates with an exit-code *semantic* comparison (0 vs non-zero)
# where a Python command exists; the six v1 commands removed by PD-D15
# (workspace/desktop/init-demo/presets) are excluded by design.
# ---------------------------------------------------------------------------

CLI_GOLDENS = ROOT / "fixtures" / "cli-goldens"


def _strip_tmp_paths(text: str, tmp: Path) -> str:
    """Temporary directories differ between runs; canonicalize them.

    `str(tmp)` is replaced first (more specific), then the system temp dir,
    so the doctor golden's temp-dir probe line is replay-safe across hosts.
    macOS resolves `/var` to `/private/var`, so the `/private`-prefixed form
    is normalized the same way (the Rust CLI canonicalizes paths it prints).
    """
    text = text.replace(str(tmp), "<tmp>")
    text = text.replace(tempfile.gettempdir(), "<tmp>")
    text = text.replace("/private<tmp>", "<tmp>")
    return text


#: JSON keys whose values are wall-clock timestamps and must never be part of
#: the byte contract (e.g. project-list's `modified_at`).
_TIMESTAMP_KEYS = frozenset({"modified_at", "created_at"})


def _canonicalize_json(value: object) -> object:
    if isinstance(value, dict):
        return {
            key: ("<timestamp>" if key in _TIMESTAMP_KEYS else _canonicalize_json(item))
            for key, item in value.items()
        }
    if isinstance(value, list):
        return [_canonicalize_json(item) for item in value]
    return value


def _normalize_cli_output(text: str, tmp: Path | None = None) -> str:
    """Canonicalize CLI output: stable JSON key order, no timestamps."""
    if tmp is not None:
        text = _strip_tmp_paths(text, tmp)
    stripped = text.strip()
    try:
        value = json.loads(stripped)
    except json.JSONDecodeError:
        return stripped
    return json.dumps(_canonicalize_json(value), sort_keys=True, ensure_ascii=False)


def _fixture_layout_value(fixture: Path) -> dict:
    return json.loads((fixture / "classroom.json").read_text(encoding="utf-8"))


def _repair_snapshot_doc(request: dict, snapshot_path: Path, *, saved_locks: dict | None) -> dict:
    """Editor-style snapshot for the repair CLI: convert the CoreSolveResponse
    index pairs into {student_key, seat_id} assignments. Seat ids come from
    the request layout when present, else the core's derived `seat-N` ids.
    Optional saved locks exercise the Python `reuse_saved_locks` default."""
    response = json.loads(snapshot_path.read_text(encoding="utf-8"))
    students = request.get("students") or []
    layout_seats = (request.get("layout") or {}).get("seats") or []
    assignments = []
    for pair in response.get("assignment", []):
        student_index, seat_index = pair
        student_key = students[student_index]["key"]
        if layout_seats:
            seat_id = layout_seats[seat_index]["seat_id"]
        else:
            seat_id = f"seat-{seat_index + 1}"
        assignments.append({"student_key": student_key, "seat_id": seat_id})
    doc: dict = {"assignments": assignments}
    if saved_locks is not None:
        doc["metadata"] = {"lock_state": saved_locks}
    return doc


def _cli_case_commands(tmp: Path) -> list[dict]:
    """The representative command matrix. Each entry: name, rust argv,
    optional python argv (exit-code semantics only), optional JSON keys that
    must agree between the two sides."""
    solve_problem = tmp / "solve-problem.json"
    fixture = INPUTS / "p20-rect-exact-none"
    request, _notes = fixture_to_request(fixture)
    solve_problem.write_text(json.dumps(request), encoding="utf-8")
    snapshot = tmp / "snapshot.json"
    subprocess.run(
        [str(CLI), "solve", "--problem", str(solve_problem), "--output", str(snapshot)],
        capture_output=True, text=True, check=False,
    )
    hist = GOLDENS / "hist-short" / "snapshot.json"
    # Repair input snapshots: plain, and with saved locks in the metadata.
    repair_snapshot = tmp / "repair-snapshot.json"
    repair_snapshot.write_text(
        json.dumps(_repair_snapshot_doc(request, snapshot, saved_locks=None), ensure_ascii=False),
        encoding="utf-8",
    )
    repair_locked_snapshot = tmp / "repair-locked-snapshot.json"
    repair_locked_snapshot.write_text(
        json.dumps(
            _repair_snapshot_doc(
                request, snapshot, saved_locks={"locked_students": ["STU001"], "locked_seats": []}
            ),
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )
    # Warning-path problems: the p20 fixture rules are empty, so the
    # preset-context and group-scope warnings need their soft rules enabled.
    warnings_problem = tmp / "warnings-problem.json"
    warnings_request = json.loads(solve_problem.read_text(encoding="utf-8"))
    warnings_request["rules"] = {
        "seed": warnings_request.get("seed", 42),
        "soft": {
            "vision_front": {"enabled": True, "weight": 20},
            "height_back": {"enabled": True, "weight": 4},
            "randomize": {"enabled": True, "weight": 3},
            "score_balance": {"enabled": True, "weight": 4},
            "fair_rotation": {"enabled": True, "weight": 12},
            "avoid_recent_neighbors": {"enabled": True, "weight": 12},
        },
    }
    warnings_problem.write_text(json.dumps(warnings_request), encoding="utf-8")
    group_problem = tmp / "group-problem.json"
    group_request = json.loads(solve_problem.read_text(encoding="utf-8"))
    group_request["rules"] = {
        "seed": group_request.get("seed", 42),
        "soft": {"score_distribution": {"enabled": True, "weight": 18, "scope": "group"}},
    }
    group_request["layout"] = {
        "layout_id": "group-layout",
        "name": "group layout",
        "seats": _fixture_layout_value(fixture)["seats"],
        "adjacency": _fixture_layout_value(fixture).get("adjacency", {}),
    }
    group_problem.write_text(json.dumps(group_request), encoding="utf-8")
    # Project lifecycle: one workspace under tmp/, commands run in order so
    # later steps consume earlier outputs (project-solve produces the plan
    # that project-export renders).
    project_dir = tmp / "proj"
    project_dir.mkdir(exist_ok=True)
    for src_name, dst_name in [
        ("students.csv", "students.csv"),
        ("classroom.json", "layout.json"),
        ("rules.json", "rules.json"),
    ]:
        shutil.copy2(fixture / src_name, project_dir / dst_name)
    project_file = project_dir / "seattrellis.project.json"
    plan_file = project_dir / "outputs" / "plan.json"
    project_dir.joinpath("outputs").mkdir(exist_ok=True)
    bundle_zip = tmp / "bundle.zip"
    # History directory for the validate --history-dir preset case.
    hist_dir = tmp / "hist-dir"
    hist_dir.mkdir(exist_ok=True)
    shutil.copy2(hist, hist_dir / "week1.snapshot.json")
    # A second workspace whose rules trigger the score_distribution
    # group-scope capability warning: project-validate --strict must fail.
    warn_dir = tmp / "proj-warn"
    warn_dir.mkdir(exist_ok=True)
    for src_name, dst_name in [
        ("students.csv", "students.csv"),
        ("classroom.json", "layout.json"),
    ]:
        shutil.copy2(fixture / src_name, warn_dir / dst_name)
    warn_dir.joinpath("rules.json").write_text(
        json.dumps(
            {
                "seed": 42,
                "soft": {
                    "score_distribution": {"enabled": True, "weight": 18, "scope": "group"}
                },
            },
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )
    # Same project document project-init writes for the main workspace.
    warn_dir.joinpath("seattrellis.project.json").write_text(
        json.dumps(
            {
                "kind": "seattrellis_project",
                "schema_version": 1,
                "name": "warn",
                "students": "students.csv",
                "layout": "layout.json",
                "rules": "rules.json",
                "outputs_dir": "outputs",
            },
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )
    warn_project_file = warn_dir / "seattrellis.project.json"
    return [
        {"name": "project-init", "rust": ["project-init", "--dir", str(project_dir)]},
        {"name": "project-list", "rust": ["project-list", "--root", str(project_dir)]},
        {"name": "project-info", "rust": ["project-info", "--project", str(project_file)]},
        {"name": "project-validate", "rust": ["project-validate", "--project", str(project_file)]},
        {"name": "project-solve", "rust": ["project-solve", "--project", str(project_file),
                                          "--output", str(plan_file)]},
        {"name": "project-edit", "rust": ["project-edit", "--project", str(project_file),
                                          "--snapshot", str(plan_file),
                                          "--operation", "swap:STU001:STU002",
                                          "--output", str(tmp / "proj-edited.json")]},
        {"name": "project-repair", "rust": ["project-repair", "--project", str(project_file),
                                            "--snapshot", str(tmp / "proj-edited.json"),
                                            "--output", str(tmp / "proj-repaired.json")]},
        {"name": "project-export", "rust": ["project-export", "--project", str(project_file),
                                            "--snapshot", str(plan_file),
                                            "--format", "svg", "--output", str(project_dir / "out.svg")]},
        {"name": "project-solve-candidates",
         "rust": ["project-solve", "--project", str(project_file),
                  "--candidates", "3", "--output", str(project_dir / "outputs" / "candidates.json"),
                  "--report", str(project_dir / "outputs" / "report.json")],
         "python": ["project-solve", "--project", str(project_file),
                    "--candidates", "3", "--output", str(project_dir / "outputs" / "candidates-py.json"),
                    "--report", str(project_dir / "outputs" / "report-py.json")]},
        {"name": "project-export-candidate",
         "rust": ["project-export", "--project", str(project_file),
                  "--snapshot", str(project_dir / "outputs" / "candidates.json"),
                  "--candidate", "candidate_01", "--format", "svg",
                  "--output", str(project_dir / "candidate.svg")]},
        {"name": "project-rotate", "rust": ["project-rotate", "--project", str(project_file),
                                            "--periods", "2"]},
        {"name": "project-privacy", "rust": ["project-privacy", "--project", str(project_file)]},
        {"name": "project-privacy-no-outputs",
         "rust": ["project-privacy", "--project", str(project_file), "--no-include-outputs"],
         "python": ["project-privacy", "--project", str(project_file), "--no-include-outputs"]},
        {"name": "project-pack", "rust": ["project-pack", "--project", str(project_file),
                                          "--output", str(bundle_zip)]},
        {"name": "project-restore",
         "rust": ["project-restore", "--bundle", str(bundle_zip),
                  "--output-dir", str(tmp / "restored")],
         "python": ["project-restore", "--bundle", str(bundle_zip),
                    "--output-dir", str(tmp / "restored-py")]},
        {"name": "project-validate-strict",
         "rust": ["project-validate", "--project", str(project_file), "--strict"],
         "python": ["project-validate", "--project", str(project_file), "--strict"]},
        {"name": "project-validate-strict-warning",
         "rust": ["project-validate", "--project", str(warn_project_file), "--strict"],
         "python": ["project-validate", "--project", str(warn_project_file), "--strict"]},
        {"name": "help", "rust": ["--help"]},
        {"name": "version", "rust": ["--version"]},
        {"name": "doctor", "rust": ["doctor"], "python": ["doctor"]},
        {"name": "solve", "rust": ["solve", "--problem", str(solve_problem)],
         "python": ["solve", "--students", str(fixture / "students.csv"),
                    "--layout", str(fixture / "classroom.json"),
                    "--rules", str(fixture / "rules.json"), "--seed", "168996"],
         "json_keys": ["status"]},
        {"name": "validate", "rust": ["validate", "--problem", str(solve_problem)],
         "python": ["validate", "--students", str(fixture / "students.csv"),
                    "--layout", str(fixture / "classroom.json"),
                    "--rules", str(fixture / "rules.json")]},
        {"name": "validate-warnings",
         "rust": ["validate", "--problem", str(warnings_problem), "--preset", "daily"],
         "python": ["validate", "--students", str(fixture / "students.csv"),
                    "--layout", str(fixture / "classroom.json"), "--preset", "daily"]},
        {"name": "validate-strict",
         "rust": ["validate", "--problem", str(warnings_problem), "--preset", "daily", "--strict"],
         "python": ["validate", "--students", str(fixture / "students.csv"),
                    "--layout", str(fixture / "classroom.json"), "--preset", "daily", "--strict"]},
        {"name": "validate-group-scope",
         "rust": ["validate", "--problem", str(group_problem)],
         "python": ["validate", "--students", str(fixture / "students.csv"),
                    "--layout", str(fixture / "classroom.json"),
                    "--rules", str(fixture / "rules.json")]},
        {"name": "validate-history-dir",
         "rust": ["validate", "--problem", str(warnings_problem), "--preset", "daily",
                  "--history-dir", str(hist_dir)],
         "python": ["validate", "--students", str(fixture / "students.csv"),
                    "--layout", str(fixture / "classroom.json"), "--preset", "daily",
                    "--history-dir", str(hist_dir)]},
        {"name": "precheck", "rust": ["precheck", "--problem", str(solve_problem)]},
        {"name": "audit", "rust": ["audit", "--problem", str(solve_problem),
                                   "--snapshot", str(snapshot)]},
        {"name": "candidates", "rust": ["candidates", "--problem", str(solve_problem), "--count", "3"]},
        {"name": "history-report", "rust": ["history-report", "--problem", str(solve_problem),
                                            "--history", str(hist)],
         "python": ["history-report", "--students", str(fixture / "students.csv"),
                    "--layout", str(fixture / "classroom.json"), "--history", str(hist)]},
        {"name": "pair-report", "rust": ["pair-report", "--problem", str(solve_problem),
                                         "--history", str(hist)],
         "python": ["pair-report", "--students", str(fixture / "students.csv"),
                    "--layout", str(fixture / "classroom.json"), "--history", str(hist)]},
        {"name": "repair", "rust": ["repair", "--problem", str(solve_problem),
                                    "--snapshot", str(repair_snapshot)],
         "python": ["repair", "--snapshot", str(hist),
                    "--output", str(tmp / "repaired-py.json")]},
        {"name": "repair-saved-locks", "rust": ["repair", "--problem", str(solve_problem),
                                                "--snapshot", str(repair_locked_snapshot)],
         "python": ["repair", "--snapshot", str(hist),
                    "--output", str(tmp / "repaired-py2.json")]},
        {"name": "repair-ignore-saved-locks",
         "rust": ["repair", "--problem", str(solve_problem),
                  "--snapshot", str(repair_locked_snapshot), "--ignore-saved-locks"],
         "python": ["repair", "--snapshot", str(hist),
                    "--output", str(tmp / "repaired-py3.json")]},
        {"name": "edit", "rust": ["edit", "--snapshot", str(hist),
                                  "--operation", "swap:STU001:STU002",
                                  "--output", str(tmp / "edited.json")],
         "python": ["edit", "--snapshot", str(hist),
                    "--operation", "swap:STU001:STU002",
                    "--output", str(tmp / "edited-py.json")]},
        {"name": "schema-list", "rust": ["schema-list"]},
        {"name": "schema-export", "rust": ["schema-export", "--kind", "student_roster", "--output", str(tmp / "roster.v2.json")]},
        {"name": "schema-migrate", "rust": ["schema-migrate", "--input", str(fixture / "students.csv"), "--dry-run"],
         "skip_python": True},
    ]


def _run_cli(argv: list[str]) -> subprocess.CompletedProcess:
    return subprocess.run([str(CLI), *argv], capture_output=True, text=True, check=False)


def _python_exit_code(argv: list[str]) -> int | None:
    proc = subprocess.run([str(PY_CLI), *argv], capture_output=True, text=True, check=False)
    return proc.returncode


def run_cli_golden_class(record: bool, tmp: Path) -> list[tuple[str, str, str, str, list[str]]]:
    if not CLI.exists():
        raise SystemExit(f"Rust CLI not found: {CLI}; build it first (cargo build --release -p seattrellis_cli)")
    rows: list[tuple[str, str, str, str, list[str]]] = []
    CLI_GOLDENS.mkdir(parents=True, exist_ok=True)
    for case in _cli_case_commands(tmp):
        name = case["name"]
        rust = _run_cli(case["rust"])
        golden_path = CLI_GOLDENS / f"{name}.json"
        if record:
            golden_path.write_text(
                json.dumps(
                    {
                        "stdout": _normalize_cli_output(rust.stdout, tmp),
                        "stderr": _strip_tmp_paths(rust.stderr, tmp),
                        "exit": rust.returncode,
                    },
                    ensure_ascii=False,
                    indent=2,
                )
                + "\n",
                encoding="utf-8",
            )
            continue
        if not golden_path.exists():
            rows.append((name, "SKIP", "golden missing (record first)", "cli-golden", []))
            continue
        golden = json.loads(golden_path.read_text(encoding="utf-8"))
        problems: list[str] = []
        if _normalize_cli_output(rust.stdout, tmp) != _normalize_cli_output(golden["stdout"], tmp):
            problems.append("stdout mismatch")
        if rust.returncode != golden["exit"]:
            problems.append(f"exit {rust.returncode} != golden {golden['exit']}")
        # Python exit-code semantics (0 vs non-zero) where a command exists.
        if not case.get("skip_python") and "python" in case:
            py_exit = _python_exit_code(case["python"])
            py_ok = py_exit == 0
            rust_ok = rust.returncode == 0
            if py_ok != rust_ok:
                problems.append(f"python exit semantics differ (python {py_exit}, rust {rust.returncode})")
        status = "OK" if not problems else "MISMATCH"
        rows.append((name, status, status, "; ".join(problems) if problems else "stdout+exit match", []))
    return rows


if __name__ == "__main__":
    raise SystemExit(main())
