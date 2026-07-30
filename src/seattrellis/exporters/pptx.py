"""Single-slide 16:9 PowerPoint export using editable native shapes."""

from __future__ import annotations

from pathlib import Path
from typing import TYPE_CHECKING, Any

from seattrellis.exporters.canvas import (
    CanvasSeat,
    SeatingCanvasDocument,
    build_seating_canvas,
    seat_font_size,
)
from seattrellis.exporters.presentation import xml_safe_text
from seattrellis.service_types import PrivacyOptions

if TYPE_CHECKING:
    from seattrellis.models.candidate import CandidatePlan
    from seattrellis.models.snapshot import SeatingSnapshot


def export_pptx(
    snapshot: "SeatingSnapshot",
    output: str | Path,
    *,
    template: str = "public",
    privacy: PrivacyOptions | None = None,
    candidate: "CandidatePlan | None" = None,
    locale: str = "zh",
) -> Path:
    """Write one widescreen slide whose seats remain individually editable."""

    document = build_seating_canvas(
        snapshot,
        template=template,
        privacy=privacy,
        candidate=candidate,
        locale=locale,
    )
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
        from pptx.enum.shapes import MSO_SHAPE  # type: ignore[import-untyped]
        from pptx.util import Inches  # type: ignore[import-untyped]
    except ImportError as exc:  # pragma: no cover - depends on optional install.
        from seattrellis.optional import MissingOptionalDependencyError

        raise MissingOptionalDependencyError("PowerPoint export", "pptx") from exc
    presentation = Presentation()
    presentation.slide_width = Inches(40 / 3)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    _add_canvas_rectangle(
        slide,
        document,
        x=0,
        y=0,
        width=document.width,
        height=document.height,
        fill=document.theme.background,
        line=document.theme.background,
        shape_kind=MSO_SHAPE.RECTANGLE,
    )
    _add_text_box(
        slide,
        document,
        document.title,
        x=80,
        y=25,
        width=document.width - 160,
        height=50,
        font_size=34,
        colour=document.theme.heading,
        bold=True,
    )
    _add_text_box(
        slide,
        document,
        document.subtitle,
        x=80,
        y=75,
        width=document.width - 160,
        height=32,
        font_size=17,
        colour=document.theme.secondary_text,
    )

    for seat in document.seats:
        _add_seat_shape(slide, document, seat, MSO_SHAPE.ROUNDED_RECTANGLE)

    _add_canvas_rectangle(
        slide,
        document,
        x=80,
        y=document.height - 56,
        width=document.width - 160,
        height=38,
        fill=document.theme.footer_fill,
        line=document.theme.footer_fill,
        shape_kind=MSO_SHAPE.ROUNDED_RECTANGLE,
    )
    _add_text_box(
        slide,
        document,
        document.footer,
        x=90,
        y=document.height - 53,
        width=document.width - 180,
        height=30,
        font_size=14,
        colour=document.theme.secondary_text,
        bold=True,
    )

    path = Path(output)
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(path))
    return path


def _add_seat_shape(
    slide: Any,
    document: SeatingCanvasDocument,
    seat: CanvasSeat,
    shape_kind: Any,
) -> None:
    from pptx.dml.color import RGBColor  # type: ignore[import-untyped]
    from pptx.enum.text import (  # type: ignore[import-untyped]
        MSO_ANCHOR,
        PP_ALIGN,
    )
    from pptx.util import Pt  # type: ignore[import-untyped]

    shape = slide.shapes.add_shape(
        shape_kind,
        _x(document, seat.x),
        _y(document, seat.y),
        _x(document, seat.width),
        _y(document, seat.height),
    )
    fill_colour = (
        document.theme.seat_fill if seat.enabled else document.theme.disabled_fill
    )
    border_colour = (
        document.theme.seat_border
        if seat.enabled
        else document.theme.disabled_border
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill_colour)
    shape.line.color.rgb = RGBColor.from_string(border_colour)
    shape.line.width = Pt(1.25)

    lines = seat.render_lines
    font_size = seat_font_size(seat) * 0.6
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.margin_left = Pt(3)
    text_frame.margin_right = Pt(3)
    text_frame.margin_top = Pt(2)
    text_frame.margin_bottom = Pt(2)
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.word_wrap = True

    for index, value in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = _safe_text(value)
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.space_after = Pt(0)
        paragraph.space_before = Pt(0)
        paragraph.line_spacing = 1
        for run in paragraph.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(max(5.0, font_size))
            run.font.bold = index == 1 or (len(lines) == 1 and index == 0)
            run.font.color.rgb = RGBColor.from_string(
                document.theme.heading
                if seat.enabled
                else document.theme.empty_text
            )


def _add_canvas_rectangle(
    slide: Any,
    document: SeatingCanvasDocument,
    *,
    x: float,
    y: float,
    width: float,
    height: float,
    fill: str,
    line: str,
    shape_kind: Any,
) -> None:
    from pptx.dml.color import RGBColor  # type: ignore[import-untyped]
    from pptx.util import Pt  # type: ignore[import-untyped]

    shape = slide.shapes.add_shape(
        shape_kind,
        _x(document, x),
        _y(document, y),
        _x(document, width),
        _y(document, height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill)
    shape.line.color.rgb = RGBColor.from_string(line)
    shape.line.width = Pt(0.1)


def _add_text_box(
    slide: Any,
    document: SeatingCanvasDocument,
    value: str,
    *,
    x: float,
    y: float,
    width: float,
    height: float,
    font_size: float,
    colour: str,
    bold: bool = False,
) -> None:
    from pptx.dml.color import RGBColor  # type: ignore[import-untyped]
    from pptx.enum.text import (  # type: ignore[import-untyped]
        MSO_ANCHOR,
        PP_ALIGN,
    )
    from pptx.util import Pt  # type: ignore[import-untyped]

    shape = slide.shapes.add_textbox(
        _x(document, x),
        _y(document, y),
        _x(document, width),
        _y(document, height),
    )
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.text = _safe_text(value)
    for run in paragraph.runs:
        run.font.name = "Aptos"
        run.font.size = Pt(font_size * 0.6)
        run.font.bold = bold
        run.font.color.rgb = RGBColor.from_string(colour)


def _x(document: SeatingCanvasDocument, value: float) -> int:
    # 40/3 inches is the PowerPoint widescreen width.
    return round(value / document.width * (40 / 3) * 914400)


def _y(document: SeatingCanvasDocument, value: float) -> int:
    return round(value / document.height * 7.5 * 914400)


def _safe_text(value: object) -> str:
    return xml_safe_text(value)
