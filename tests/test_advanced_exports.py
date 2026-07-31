from __future__ import annotations

from datetime import datetime, timezone
from xml.etree import ElementTree
from zipfile import ZipFile

import pytest

from seattrellis.exporters import build_seating_canvas, export_snapshot
from seattrellis.exporters.canvas import CanvasTheme
from seattrellis.exporters.pptx import _safe_text as pptx_safe_text
from seattrellis.io.json_files import write_json_model
from seattrellis.models.candidate import CandidatePlan, CandidateSet
from seattrellis.models.layout import ClassroomLayout, SeatNode
from seattrellis.models.rules import RuleSet
from seattrellis.models.snapshot import SeatAssignment, SeatingSnapshot
from seattrellis.models.student import Student
from seattrellis.scoring import score_snapshot
from seattrellis.service import export as service_export
from seattrellis.service_types import ExportRequest, PageOptions, PrivacyOptions


FIXED_CREATED_AT = datetime(2026, 7, 30, 8, 9, 10, tzinfo=timezone.utc)


def _sensitive_snapshot() -> SeatingSnapshot:
    student = Student(
        student_id="S1",
        name='Ada\x01</text><script>alert("seat")</script>',
        score=88,
        height_cm=172,
        vision="SECRET_VISION",
        notes="SECRET_NOTE<&",
        needs=["SECRET_NEED"],
        tags=["SECRET_TAG"],
    )
    layout = ClassroomLayout(
        name="<Unsafe & Classroom>",
        seats=[
            SeatNode(seat_id='S<1>"', row=1, col=1),
            SeatNode(seat_id="S2", row=1, col=3, enabled=False),
            SeatNode(seat_id="S3", row=2, col=2),
        ],
    )
    return SeatingSnapshot(
        created_at=FIXED_CREATED_AT,
        students=[student],
        layout=layout,
        rules=RuleSet(),
        assignments=[
            SeatAssignment(
                student_key=student.key,
                student_name=student.display_name,
                seat_id='S<1>"',
            )
        ],
        solver_status="FEASIBLE",
    )


def _candidate(snapshot: SeatingSnapshot) -> CandidatePlan:
    return CandidatePlan(
        candidate_id="candidate_<01>",
        snapshot=snapshot,
        score=score_snapshot(snapshot),
        hard_constraints_satisfied=True,
    )


def _all_xml_text(root: ElementTree.Element) -> str:
    return "".join(value for value in root.itertext() if value)


def test_canvas_has_stable_widescreen_geometry_and_shared_privacy_rules() -> None:
    snapshot = _sensitive_snapshot()

    public = build_seating_canvas(snapshot, template="public", locale="zh")
    teacher = build_seating_canvas(snapshot, template="teacher", locale="en")
    teacher_again = build_seating_canvas(snapshot, template="teacher", locale="en")

    assert public.width == 1600
    assert public.height == 900
    assert public.aspect_ratio == pytest.approx(16 / 9)
    assert teacher == teacher_again
    assert [seat.seat_id for seat in teacher.seats] == ['S<1>"', "S2", "S3"]
    assert [(seat.x, seat.y, seat.width, seat.height) for seat in public.seats] == [
        (seat.x, seat.y, seat.width, seat.height) for seat in teacher.seats
    ]

    public_text = " ".join(
        line for seat in public.seats for line in seat.text_lines
    )
    teacher_text = " ".join(
        line for seat in teacher.seats for line in seat.text_lines
    )
    assert "SECRET_NOTE" not in public_text
    assert "SECRET_NEED" not in public_text
    assert "88.0" not in public_text
    assert "Student ID: S1" in teacher_text
    assert "Score: 88.0" in teacher_text
    assert "Vision: SECRET_VISION" in teacher_text
    assert "Special needs: SECRET_NEED, SECRET_TAG" in teacher_text
    assert "Notes: SECRET_NOTE<&" in teacher_text


def test_canvas_rejects_renderer_attribute_injection() -> None:
    with pytest.raises(ValueError, match="six-digit hexadecimal colour"):
        CanvasTheme(background='FFFFFF" onload="alert(1)')

    assert pptx_safe_text("safe\x01text") == "safe\ufffdtext"


def test_svg_is_safe_self_contained_and_uses_only_native_elements(tmp_path) -> None:
    snapshot = _sensitive_snapshot()
    request = ExportRequest(
        output_format="svg",
        output_path=tmp_path / "public.svg",
        template="public",
        locale="en",
    )

    output = export_snapshot(snapshot, request=request)
    raw = output.read_text(encoding="utf-8")
    root = ElementTree.fromstring(raw)
    local_tags = {element.tag.rsplit("}", 1)[-1] for element in root.iter()}

    assert local_tags <= {"svg", "rect", "text", "tspan"}
    assert "<script" not in raw.lower()
    assert "<!doctype" not in raw.lower()
    assert "foreignobject" not in raw.lower()
    assert "href=" not in raw.lower()
    assert "SECRET_NOTE" not in raw
    assert "SECRET_NEED" not in raw
    assert "SECRET_VISION" not in raw
    assert "88.0" not in raw
    assert "\ufffd" in _all_xml_text(root)
    assert '</text><script>alert("seat")</script>' in _all_xml_text(root)


def test_service_export_uses_the_requested_candidate_for_svg(tmp_path) -> None:
    recommended_snapshot = _sensitive_snapshot()
    selected_snapshot = recommended_snapshot.model_copy(deep=True)
    selected_snapshot.layout.name = "Selected Classroom"
    selected_snapshot.assignments[0].student_name = "Selected Student"
    recommended = CandidatePlan(
        candidate_id="candidate_recommended",
        snapshot=recommended_snapshot,
        score=score_snapshot(recommended_snapshot),
        hard_constraints_satisfied=True,
    )
    selected = CandidatePlan(
        candidate_id="candidate_selected",
        snapshot=selected_snapshot,
        score=score_snapshot(selected_snapshot),
        hard_constraints_satisfied=True,
    )
    artifact = write_json_model(
        CandidateSet(
            candidates=[recommended, selected],
            recommended_candidate_id=recommended.candidate_id,
        ),
        tmp_path / "candidate-selection.json",
    )

    output = service_export(
        snapshot_path=artifact,
        request=ExportRequest(
            output_format="svg",
            output_path=tmp_path / "selected.svg",
            template="report",
            candidate_id=selected.candidate_id,
        ),
    ).read_text(encoding="utf-8")

    assert "Selected Classroom" in output
    assert "Selected Student" in output
    assert "candidate_selected" in output
    assert "candidate_recommended" not in output


def test_svg_teacher_fields_are_localized_and_anonymization_is_consistent(
    tmp_path,
) -> None:
    snapshot = _sensitive_snapshot()
    teacher = export_snapshot(
        snapshot,
        request=ExportRequest(
            output_format="svg",
            output_path=tmp_path / "teacher.svg",
            template="teacher",
            locale="zh",
        ),
    ).read_text(encoding="utf-8")
    anonymous = export_snapshot(
        snapshot,
        request=ExportRequest(
            output_format="svg",
            output_path=tmp_path / "anonymous.svg",
            template="teacher",
            locale="en",
            privacy=PrivacyOptions(
                hide_scores=True,
                hide_notes=True,
                hide_special_needs=True,
                anonymize=True,
                show_height=False,
                show_vision=False,
            ),
        ),
    ).read_text(encoding="utf-8")

    assert "学号: S1" in teacher
    assert "成绩: 88.0" in teacher
    assert "备注: SECRET_NOTE&lt;&amp;" in teacher
    assert "Student 01" in anonymous
    for secret in [
        "S1",
        "Ada",
        "SECRET_NOTE",
        "SECRET_NEED",
        "SECRET_TAG",
        "SECRET_VISION",
        "88.0",
        "172.0",
    ]:
        assert secret not in anonymous


@pytest.mark.parametrize("output_format", ["svg", "pptx"])
def test_advanced_exports_reject_all_candidate_scope_explicitly(
    tmp_path,
    output_format,
) -> None:
    snapshot = _sensitive_snapshot()
    candidate = _candidate(snapshot)
    artifact = write_json_model(
        CandidateSet(
            candidates=[candidate],
            recommended_candidate_id=candidate.candidate_id,
        ),
        tmp_path / "candidates.json",
    )
    request = ExportRequest(
        output_format=output_format,
        output_path=tmp_path / f"all.{output_format}",
        candidate_scope="all",
    )

    with pytest.raises(
        ValueError,
        match="does not support candidate_scope='all'; select one candidate",
    ):
        service_export(snapshot_path=artifact, request=request)


@pytest.mark.parametrize("output_format", ["svg", "pptx"])
def test_advanced_exports_use_fixed_canvas_and_require_report_candidate(
    tmp_path,
    output_format,
) -> None:
    snapshot = _sensitive_snapshot()
    default_request = ExportRequest(output_format=output_format)

    assert default_request.page == PageOptions()
    assert default_request.resolved_output_path.suffix == f".{output_format}"

    with pytest.raises(ValueError, match="fixed 16:9 canvas"):
        export_snapshot(
            snapshot,
            request=ExportRequest(
                output_format=output_format,
                output_path=tmp_path / f"landscape.{output_format}",
                page=PageOptions(orientation="landscape"),
            ),
        )
    with pytest.raises(ValueError, match="requires a candidate plan"):
        export_snapshot(
            snapshot,
            request=ExportRequest(
                output_format=output_format,
                output_path=tmp_path / f"report.{output_format}",
                template="report",
            ),
        )


def test_pptx_is_widescreen_and_uses_editable_native_shapes(tmp_path) -> None:
    pytest.importorskip("pptx")
    from pptx import Presentation

    snapshot = _sensitive_snapshot()
    output = export_snapshot(
        snapshot,
        request=ExportRequest(
            output_format="pptx",
            output_path=tmp_path / "teacher.pptx",
            template="teacher",
            locale="en",
        ),
    )

    presentation = Presentation(output)
    assert len(presentation.slides) == 1
    assert presentation.slide_width * 9 == presentation.slide_height * 16
    assert len(presentation.slides[0].shapes) >= len(snapshot.layout.seats) + 5

    with ZipFile(output) as archive:
        names = archive.namelist()
        slide_xml = archive.read("ppt/slides/slide1.xml").decode("utf-8")
        relationship_xml = "\n".join(
            archive.read(name).decode("utf-8")
            for name in names
            if name.endswith(".rels")
        )
    root = ElementTree.fromstring(slide_xml)
    text = _all_xml_text(root)

    assert not any(name.startswith("ppt/media/") for name in names)
    assert not any(name.endswith("vbaProject.bin") for name in names)
    assert 'TargetMode="External"' not in relationship_xml
    assert "<p:pic" not in slide_xml
    assert "<script" not in slide_xml.lower()
    assert "Student ID: S1" in text
    assert "Score: 88.0" in text
    assert "Notes: SECRET_NOTE<&" in text
    assert '</text><script>alert("seat")</script>' in text
